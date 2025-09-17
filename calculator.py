import os
import webbrowser
from yt_dlp import YoutubeDL
import psutil
import tkinter as tk
from tkinter import ttk
import threading
import subprocess
import time
from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from pptx import Presentation
from pptx.util import Inches as PPTInches
import requests_cache
import pandas as pd
from openmeteo_requests import Client
from requests.adapters import HTTPAdapter, Retry
import vlc
import tkinter as tk
from tkinter import filedialog
import time
# ---------------- WELCOME ----------------
def welcome():
    print(r"""
  ____  _            _       _             
 / ___|| | ___   ___| | __  | | _____ _ __ 
 \___ \| |/ _ \ / __| |/ /  | |/ / _ \ '__|
  ___) | | (_) | (__|   < |_|   <  __/ |   
 |____/|_|\___/ \___|_|\_\(_)_|\_\___|_|   

Welcome to the Super Calculator!
""")

def exit_program():
    print("Exiting the program. Goodbye!")

# ---------------- BASIC MATH ----------------
def math():
    retry = "y"
    while retry.lower() == "y":
        op = input("Enter operation (add, sub, mul, div, pow, sqr, sqrt): ").lower()
        
        if op in ["add", "sub", "mul", "div", "pow"]:
            a = float(input("Enter first number: "))
            b = float(input("Enter second number: "))
        
        if op == "add":
            print("Addtion really easy huh? hard for ya NOOB!")
            print("The sum is:", a+b)
        elif op == "sub":
            print("Bro why the hell you need a calculator to subtract??? you a kid?")
            print("The difference is:", a-b)
        elif op == "mul":
            print("Multiplication huh? you know times tables right?")
            print("The product is:", a*b)
        elif op == "div":
            if b != 0:
                print("Division huh? you really need a calculator for that???")
                print("The quotient is:", a/b)
            else:
                print("NOOB! how could you divide by zero???")
                print("Cannot divide by zero!")
        elif op == "pow":
            print("Exponentiation huh? bro you know what that means???")
            print("The result is:", a**b)
        elif op == "sqr":
            print("Square huh? what will you do with that???")
            a = float(input("Enter number: "))
            print("Square is:", a**2)
        elif op == "sqrt":
            print("Square root huh? yeah sure bro why not but still nOOB!")
            a = float(input("Enter number: "))
            print("Square root is:", a**0.5)
        else:
            print("bro are you intoxicated or dumb???")
            print("Invalid operation!")
        
        retry = input("Do you want to retry? noob! (y/n): ")

# ---------------- VOLUMES ----------------
def volumes():
        retry = "y"
        while retry.lower() == "y":
            shape = input("Enter shape (cube, cuboid, cylinder, cone, sphere): ").lower()
            
            if shape == "cube":
                print("Cube huh? bro just cube (^3) its sides???")
                a = float(input("Enter side length: "))
                print("Volume is:", a**3)
            elif shape == "cuboid":
                print("Cuboid huh? bro just multiply its length, breadth and height???")
                l = float(input("Length: "))
                b = float(input("Breadth: "))
                h = float(input("Height: "))
                print("Volume is:", l*b*h)
            elif shape == "cylinder":
                print("Cylinder huh? lol you are noob! its pi*r^2*h")
                r = float(input("Radius: "))
                h = float(input("Height: "))
                print("Volume is:", 3.14*r*r*h)
            elif shape == "cone":
                print("Cone huh? hmm its your icecream right??? ")
                r = float(input("Radius: "))
                h = float(input("Height: "))
                print("Volume is:", (1/3)*3.14*r*r*h)
            elif shape == "sphere":
                print("Sphere huh? its your head right??? ")
                r = float(input("Radius: "))
                print("Volume is:", (4/3)*3.14*r**3)
            else:
                print("bro is that your shape??")
                print("Invalid shape!")        
                retry = input("Do you want to retry? (y/n): ")

# ---------------- SURFACE AREAS ----------------
def surface_areas():
    retry = "y"
    while retry.lower() == "y":
        shape = input("Enter shape (cube, cuboid, cylinder, cone, sphere): ").lower()
        
        if shape == "cube":
            print("bro you may better consult your geometry book for this")
            a = float(input("Enter side length: "))
            print("Surface area is:", 6*a*a)
        elif shape == "cuboid":
            print("bro just multiply 2 with sum of all sides")
            l = float(input("Length: "))
            b = float(input("Breadth: "))
            h = float(input("Height: "))
            print("Surface area is:", 2*(l*b + b*h + h*l))
        elif shape == "cylinder":
            print("bro you are really dumb if you need a calculator for this just use formula 2πr(r+h)")
            r = float(input("Radius: "))
            h = float(input("Height: "))
            print("Surface area is:", 2*3.14*r*(r+h))
        elif shape == "cone":
            r = float(input("Radius: "))
            h = float(input("Height: "))
            l = (r*r + h*h)**0.5
            print("Surface area is:", 3.14*r*(r+l))
        elif shape == "sphere":
            r = float(input("Radius: "))
            print("Surface area is:", 4*3.14*r*r)
        else:
            print("Invalid shape!")
        
        retry = input("Do you want to retry? (y/n): ")

# ---------------- UNIT CONVERTER ----------------
def coverter():
    retry = "y"
    while retry.lower() == "y":
        conv = input("Enter conversion (cm to m, m to cm, km to m, m to km): ").lower()
        
        if conv == "cm to m":
            cm = float(input("Enter centimeters: "))
            print("Meters:", cm / 100)
        elif conv == "m to cm":
            m = float(input("Enter meters: "))
            print("Centimeters:", m * 100)
        elif conv == "km to m":
            km = float(input("Enter kilometers: "))
            print("Meters:", km * 1000)
        elif conv == "m to km":
            m = float(input("Enter meters: "))
            print("Kilometers:", m / 1000)
        else:
            print("Invalid conversion!")
        
        retry = input("Do you want to retry? (y/n): ")

# ---------------- PROGRAMMER CALC ----------------
def programmer_calc():
    retry = "y"
    while retry.lower() == "y":
        op = input("Enter operation (and, or, xor, not, lshift, rshift): ").lower()
        
        if op in ["and", "or", "xor", "lshift", "rshift"]:
            a = int(input("Enter first integer: "))
            b = int(input("Enter second integer: "))
        
        if op == "and":
            print("Result is:", a & b)
        elif op == "or":
            print("Result is:", a | b)
        elif op == "xor":
            print("Result is:", a ^ b)
        elif op == "not":
            a = int(input("Enter integer: "))
            print("Result is:", ~a)
        elif op == "lshift":
            print("Result is:", a << b)
        elif op == "rshift":
            print("Result is:", a >> b)
        else:
            print("Invalid operation!")
        
        retry = input("Do you want to retry? (y/n): ")

# ---------------- BIN/DEC/HEX/OCT ----------------
def bin_decimal_hex_oct():
    retry = "y"
    while retry.lower() == "y":
        num = int(input("Enter an integer: "))
        print("Binary:", bin(num))
        print("Decimal:", num)
        print("Hexadecimal:", hex(num))
        print("Octal:", oct(num))
        
        retry = input("Do you want to retry? (y/n): ")

# ---------------- DATA SIZE ----------------
def data():
    retry = "y"
    while retry.lower() == "y":
        size = float(input("Enter size: "))
        unit = input("Enter unit (B, KB, MB, GB, TB): ").upper()
        
        if unit == "B":
            print(f"{size} B = {size/1024} KB = {size/(1024**2)} MB = {size/(1024**3)} GB = {size/(1024**4)} TB")
        elif unit == "KB":
            print(f"{size} KB = {size*1024} B = {size/1024} MB = {size/(1024**2)} GB = {size/(1024**3)} TB")
        elif unit == "MB":
            print(f"{size} MB = {size*(1024**2)} B = {size*1024} KB = {size/1024} GB = {size/(1024**2)} TB")
        elif unit == "GB":
            print(f"{size} GB = {size*(1024**3)} B = {size*(1024**2)} KB = {size*1024} MB = {size/1024} TB")
        elif unit == "TB":
            print(f"{size} TB = {size*(1024**4)} B = {size*(1024**3)} KB = {size*(1024**2)} MB = {size*1024} GB")
        else:
            print("Invalid unit!")
        
        retry = input("Do you want to retry? (y/n): ")

# ---------------- TEMPERATURE ----------------
def temprature():
    retry = "y"
    while retry.lower() == "y":
        temp = float(input("Enter temperature: "))
        unit = input("Enter unit (C, F, K): ").upper()
        
        if unit == "C":
            print(f"{temp} C = {(temp * 9/5) + 32} F = {temp + 273.15} K")
        elif unit == "F":
            print(f"{temp} F = {(temp - 32) * 5/9} C = {(temp - 32) * 5/9 + 273.15} K")
        elif unit == "K":
            print(f"{temp} K = {temp - 273.15} C = {(temp - 273.15) * 9/5 + 32} F")
        else:
            print("Invalid unit!")
        
        retry = input("Do you want to retry? (y/n): ")

# ---------------- ENERGY ----------------
def energy():
    retry = "y"
    while retry.lower() == "y":
        a = input("Enter joules or kcal? (j/k): ").lower()
        if a == "j":
            joules = float(input("Enter energy in Joules: "))
            print(f"{joules} J = {joules/4184} kcal = {joules/1000} kJ")
        elif a == "k":
            kcal = float(input("Enter energy in kcal: "))
            print(f"{kcal} kcal = {kcal*4184} J = {kcal*4.184} kJ")
        else:
            print("Invalid input!")
        
        retry = input("Do you want to retry? (y/n): ")

# ---------------- USER SLEEP TIMER ----------------
def sleep_timer():
    import time as t_module  # safe import, avoids conflicts

    retry = "y"
    while retry.lower() == "y":
        user_input = input("Enter sleep time (hh:mm:ss): ").strip()
        try:
            h, m, s = map(int, user_input.split(":"))
            total_seconds = h*3600 + m*60 + s
        except Exception:
            print("Invalid format! Use hh:mm:ss")
            continue

        print(f"Sleeping for {h}h {m}m {s}s...")
        start_time = t_module.time()
        end_time = start_time + total_seconds

        try:
            while t_module.time() < end_time:
                remaining = int(end_time - t_module.time())
                rh = remaining // 3600
                rm = (remaining % 3600) // 60
                rs = remaining % 60
                print(f"Time left: {rh:02d}:{rm:02d}:{rs:02d}", end="\r")
                t_module.sleep(1)
            print("\nTime's up! ✅")
        except KeyboardInterrupt:
            print("\nTimer cancelled by user ❌")

        retry = input("Do you want to set another timer? (y/n): ")

# ---------------- BROWSE ----------------
def browse():
    a = 'y'
    while a.lower() == 'y':
        b = input("Enter website URL: ").strip()
        if not b.startswith(("http://", "https://")):
            b = "https://" + b
        try:
            webbrowser.open(b)
            print(f"Opening: {b}")
        except Exception as e:
            print(f"Error: {e}")
        
        a = input("Do you want to open another website? (y/n): ").strip()

# ---------------- YT-DLP ----------------
def download_video():
    url = input("Enter YouTube URL: ").strip()
    ydl_opts = {
        'format': 'best',
        'outtmpl': '%(title)s.%(ext)s',  # saves as video_title.mp4
    }
    try:
        with YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
        print("Download complete ✅")
    except Exception as e:
        print(f"Error: {e}")

# ---------------- FILE MANAGER + SHREDDER ----------------
def secure_delete(filename, passes=3):
    try:
        if not os.path.isfile(filename):
            print("File not found!")
            return

        filesize = os.path.getsize(filename)
        with open(filename, "ba+", buffering=0) as f:
            for i in range(passes):
                f.seek(0)
                if i % 2 == 0:
                    f.write(b"\x00" * filesize)  # overwrite with 0s
                else:
                    f.write(b"\xFF" * filesize)  # overwrite with 1s
                f.flush()
                os.fsync(f.fileno())
                print(f"Pass {i+1}/{passes} complete ✅")

        os.remove(filename)
        print(f"Securely deleted '{filename}' 🔒")
    except Exception as e:
        print(f"Error: {e}")

def file_manager():
    retry = "y"
    while retry.lower() == "y":
        print("\n--- File Manager ---")
        print("1. List files in current directory")
        print("2. Delete a file (normal)")
        print("3. Change directory")
        print("4. Secure Delete (overwrite with 0s/1s)")
        print("5. Exit File Manager")

        choice = input("Enter your choice: ")

        if choice == "1":
            files = os.listdir(".")
            print("\nFiles in current directory:")
            for f in files:
                print(" -", f)

        elif choice == "2":
            filename = input("Enter file name to delete: ").strip()
            if os.path.exists(filename):
                confirm = input(f"Are you sure you want to delete '{filename}'? (y/n): ").lower()
                if confirm == "y":
                    os.remove(filename)
                    print(f"Deleted '{filename}' ✅")
                else:
                    print("Cancelled ❌")
            else:
                print("File not found!")

        elif choice == "3":
            path = input("Enter new directory path: ").strip()
            try:
                os.chdir(path)
                print(f"Changed directory to {os.getcwd()}")
            except Exception as e:
                print(f"Error: {e}")

        elif choice == "4":
            filename = input("Enter file name to securely delete: ").strip()
            if os.path.exists(filename):
                confirm = input(f"⚠️ SECURE DELETE! Overwrite & remove '{filename}'? (y/n): ").lower()
                if confirm == "y":
                    secure_delete(filename)
                else:
                    print("Cancelled ❌")
            else:
                print("File not found!")

        elif choice == "5":
            break
        else:
            print("Invalid choice!")

        retry = input("Do you want to continue in File Manager? (y/n): ")

# ---------------- SYSTEM MONITOR ----------------
def system_monitor():
    print("\nLaunching CPU & RAM Monitor... (Close window to return to menu)")
    
    root = tk.Tk()
    root.title("CPU & RAM Monitor")
    root.geometry("450x250")

    # CPU
    tk.Label(root, text="CPU Usage (%)", font=("Arial", 12)).pack(pady=10)
    cpu_bar = ttk.Progressbar(root, orient='horizontal', length=400, mode='determinate')
    cpu_bar.pack(pady=5)

    # RAM
    tk.Label(root, text="RAM Usage (%)", font=("Arial", 12)).pack(pady=10)
    ram_bar = ttk.Progressbar(root, orient='horizontal', length=400, mode='determinate')
    ram_bar.pack(pady=5)

    # Update function
    def update_bars():
        while True:
            cpu_bar['value'] = psutil.cpu_percent(interval=1)
            ram_bar['value'] = psutil.virtual_memory().percent
            root.update_idletasks()

    t = threading.Thread(target=update_bars, daemon=True)
    t.start()

    root.mainloop()
def word_gui():
    def save_docx():
        filename = filedialog.asksaveasfilename(defaultextension=".docx")
        if filename:
            doc = Document()
            doc.add_heading(title_entry.get(), 0)
            doc.add_paragraph(text_area.get("1.0", tk.END))
            if img_path.get():
                doc.add_picture(img_path.get(), width=Inches(2))
            doc.save(filename)
            messagebox.showinfo("Saved", f"Document saved as {filename}")

    def open_image():
        path = filedialog.askopenfilename(filetypes=[("Image Files", "*.png;*.jpg;*.jpeg;*.bmp")])
        if path:
            img_path.set(path)
            img_label.config(text=f"Image: {path}")

    win = tk.Tk()
    win.title("Word-like Editor")
    win.geometry("700x500")
    tk.Label(win, text="Title:").pack()
    title_entry = tk.Entry(win, font=("Arial", 14))
    title_entry.pack(fill="x")
    tk.Label(win, text="Content:").pack()
    text_area = tk.Text(win, wrap="word", font=("Arial", 12))
    text_area.pack(expand=True, fill="both")
    img_path = tk.StringVar()
    img_label = tk.Label(win, text="No image selected")
    img_label.pack()
    tk.Button(win, text="Insert Image", command=open_image).pack()
    tk.Button(win, text="Save as DOCX", command=save_docx).pack(pady=10)
    win.mainloop()

# ---------------- EXCEL-LIKE EDITOR (GUI) ----------------
def excel_gui():
    def save_xlsx():
        filename = filedialog.asksaveasfilename(defaultextension=".xlsx")
        if filename:
            wb = Workbook()
            ws = wb.active
            for r in range(rows):
                for c in range(cols):
                    ws.cell(row=r+1, column=c+1, value=entries[r][c].get())
            if img_path.get():
                try:
                    img = XLImage(img_path.get())
                    ws.add_image(img, "A10")
                except Exception as e:
                    messagebox.showerror("Image Error", str(e))
            wb.save(filename)
            messagebox.showinfo("Saved", f"Excel file saved as {filename}")

    def open_image():
        path = filedialog.askopenfilename(filetypes=[("Image Files", "*.png;*.jpg;*.jpeg;*.bmp")])
        if path:
            img_path.set(path)
            img_label.config(text=f"Image: {path}")

    win = tk.Tk()
    win.title("Excel-like Editor")
    win.geometry("800x400")
    rows, cols = 10, 5
    entries = [[tk.Entry(win, width=15) for _ in range(cols)] for _ in range(rows)]
    for r, row in enumerate(entries):
        for c, entry in enumerate(row):
            entry.grid(row=r, column=c)
    img_path = tk.StringVar()
    img_label = tk.Label(win, text="No image selected")
    img_label.grid(row=rows, column=0, columnspan=cols)
    tk.Button(win, text="Insert Image", command=open_image).grid(row=rows+1, column=0, columnspan=cols)
    tk.Button(win, text="Save as XLSX", command=save_xlsx).grid(row=rows+2, column=0, columnspan=cols)
    win.mainloop()

# ---------------- POWERPOINT-LIKE EDITOR (GUI) ----------------
def pptx_gui():
    def save_pptx():
        filename = filedialog.asksaveasfilename(defaultextension=".pptx")
        if filename:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title_entry.get()
            slide.placeholders[1].text = content_text.get("1.0", tk.END)
            if img_path.get():
                slide.shapes.add_picture(img_path.get(), PPTInches(2), PPTInches(2), width=PPTInches(2))
            prs.save(filename)
            messagebox.showinfo("Saved", f"PowerPoint file saved as {filename}")

    def open_image():
        path = filedialog.askopenfilename(filetypes=[("Image Files", "*.png;*.jpg;*.jpeg;*.bmp")])
        if path:
            img_path.set(path)
            img_label.config(text=f"Image: {path}")

    win = tk.Tk()
    win.title("PowerPoint-like Editor")
    win.geometry("600x400")
    tk.Label(win, text="Slide Title:").pack()
    title_entry = tk.Entry(win, width=40)
    title_entry.pack()
    tk.Label(win, text="Slide Content:").pack()
    content_text = tk.Text(win, height=10, width=50)
    content_text.pack()
    img_path = tk.StringVar()
    img_label = tk.Label(win, text="No image selected")
    img_label.pack()
    tk.Button(win, text="Insert Image", command=open_image).pack()
    tk.Button(win, text="Save as PPTX", command=save_pptx).pack(pady=10)
import subprocess

def launch_game():
    retry = 'y'
    while retry.lower() == 'y':
        print("\nWelcome to Game Zone")
        print("Select game from the following:")
        print("1. Minecraft\n2. Counter-Strike\n3. Valorant")

        try:
            a = int(input("Enter your choice: "))
        except ValueError:
            print("Enter a valid number!")
            continue

        if a == 1:
            game_path = r"C:\Users\tejas\AppData\Roaming\.minecraft\TLauncher.exe"
        elif a == 2:
            game_path = r"C:\Games\Counter-Strike Standard Edition\cstike.exe"
        elif a == 3:
            game_path = r"D:\Riot Games\Riot Client\RiotClientServices.exe"
        else:
            print("Invalid choice!")
            retry = input("Do you want to retry? (y/n): ")
            continue

        if os.path.exists(game_path):
            try:
                subprocess.Popen(game_path)
                print(f"Launching {game_path}...")
            except Exception as e:
                print(f"Error launching game: {e}")
        else:
            print("Game not found!")

        retry = input("Do you want to launch another game? (y/n): ")
def weather_data(la=28.6139, lo=77.2090):  # Default: New Delhi coordinates
    """
    Fetch hourly temperature data from Open-Meteo and return as pandas DataFrame.
    la: latitude
    lo: longitude
    """
    try:
        # ---------------- Setup session with cache and retry ----------------
        session = requests_cache.CachedSession('.cache', expire_after=3600)
        retries = Retry(total=5, backoff_factor=0.2, status_forcelist=[500, 502, 503, 504])
        adapter = HTTPAdapter(max_retries=retries)
        session.mount("http://", adapter)
        session.mount("https://", adapter)

        client = Client(session=session)

        # ---------------- API Request ----------------
        url = "https://api.open-meteo.com/v1/forecast"
        params = {
            "latitude": la,
            "longitude": lo,
            "hourly": "temperature_2m",
            "forecast_days": 16,
        }

        responses = client.weather_api(url, params=params)
        if not responses:
            print("No response from Open-Meteo API.")
            return None

        response = responses[0]

        print(f"Coordinates: {response.Latitude()}°N {response.Longitude()}°E")
        print(f"Elevation: {response.Elevation()} m asl")
        print(f"Timezone offset (seconds from GMT+0): {response.UtcOffsetSeconds()}s")

        # ---------------- Process hourly data ----------------
        hourly = response.Hourly()
        temperature_values = hourly.Variables(0).ValuesAsNumpy()
        times_start = pd.to_datetime(hourly.Time(), unit='s', utc=True)
        times_end = pd.to_datetime(hourly.TimeEnd(), unit='s', utc=True)
        freq_seconds = hourly.Interval()
        hourly_data = {
            "date": pd.date_range(start=times_start, end=times_end, freq=pd.Timedelta(seconds=freq_seconds), inclusive="left"),
            "temperature_2m": temperature_values
        }
        hourly_df = pd.DataFrame(hourly_data)
        print("\nHourly Temperature Data (°C):\n")
        print(hourly_df.head(24)) 
        return hourly_df
    except Exception as e:
        print(f"Error fetching weather data: {e}")
        return None
def startweather():
    a=float(input("Enter latitude:"))
    b=float(input("Enter longitude:"))
    weather_data(a,b)
def video_player():
    root = tk.Tk()
    root.withdraw()
    filepath = filedialog.askopenfilename(
        title="Select a video file",
        filetypes=[("Video Files", "*.mp4;*.avi;*.mov;*.mkv")]
    )
    if not filepath:
        print("No file selected. Returning to menu...")
        return
    print(f"Playing video with audio: {filepath}")
    print("Controls: close VLC window to stop playback.")
    player = vlc.MediaPlayer(filepath)
    player.play()
    while True:
        state = player.get_state()
        if state in [vlc.State.Ended, vlc.State.Error]:
            break
        time.sleep(0.5)

    print("Playback finished.")
# ---------------- MENU ----------------
def menu():
    welcome()
    while True:
        print("\nSelect an option:")
        print("1. Basic Math Operations")
        print("2. Volume Calculations")
        print("3. Surface Area Calculations")
        print("4. Unit Converter (length)")
        print("5. Programmer Calculator (bitwise ops)")
        print("6. Number System Converter (bin/dec/hex/oct)")
        print("7. Data Size Converter (B/KB/MB/GB/TB)")
        print("8. Temperature Converter (C/F/K)")
        print("9. Energy Converter (J/kcal/kJ)")
        print("10. Sleep Timer (hh:mm:ss)")
        print("11. Open Website (Browser)")
        print("12. Download YouTube Video")
        print("13. File Manager + Shredder")
        print("14. CPU & RAM Monitor")
        print("15. Word-like Editor (GUI)")
        print("16. Excel-like Editor (GUI)")
        print("17. PowerPoint-like Editor (GUI)")
        print("18. Launch Game")
        print("19. Weather Data Fetcher")
        print("20. Video Player")
        print("21. Exit")
        choice = input("Enter your choice (1-21): ").lower()
        
        if choice in ["1", "m", "math"]:
            print("Oh so lets help your dumb brain to do some basic calculations")
            math()
        elif choice in ["2", "v", "volume"]:
            print("I think you skipped your middle school NOOB!")
            volumes()
        elif choice in ["3", "s", "surface"]:
            print("I think you skipped your middle school NOOB!")
            surface_areas()
        elif choice in ["4", "u", "unit"]:
            print("ooh so convert units huh? really bro you are lazy!")
            coverter()
        elif choice in ["5", "p", "programmer"]:
            print("You must be a programmer if you are using this")
            programmer_calc()
        elif choice in ["6", "n", "number"]:
            print("Bro you really a programmer or just wasting your battery juice???")
            bin_decimal_hex_oct()
        elif choice in ["7", "d", "data"]:
            print("Why you need to convert data sizes???")
            data()
        elif choice in ["8", "t", "temp", "temperature"]:
            print("Bro why you need to convert temprature???")
            temprature()
        elif choice in ["9", "e", "energy"]:
            print("Bro why you need to convert energy??? Ah lazy boy huh?")
            energy()
        elif choice in ["10", "ti", "time"]:
            print("Oh lazy guy need a sleep go on set a timer! I am already tired of you")
            sleep_timer()
        elif choice in ["11", "b", "browse"]:
            print("Oh so you need to open a website huh? you know how to type a url right?")
            browse()
        elif choice in ["12", "yt", "youtube"]:
            print("Downloading youtube videos huh? I hope you are not downloading copyrighted content")
            download_video()
        elif choice in ["13", "fm", "file"]:
            print("File manager + Shredder (really dont use your system files) !")
            file_manager()
        elif choice in ["14", "cpu", "ram", "monitor"]:
            print("Ah checking your potato pc huh? I know the feeling ; _ ;")
            system_monitor()
        elif choice in ["15", "word"]:
            print("bro really you need a word processor???")
            word_gui()
        elif choice in ["16", "excel"]:
            print("bro really you need a excel processor???")
            excel_gui()
        elif choice in ["17", "pptx", "powerpoint"]:
            print("bro really you need a powerpoint processor???")
            pptx_gui()
        elif choice in ["18", "game", "games"]:
            print("Hey see that lazy guys fingure out how to launch games")
            launch_game()
        elif choice in ["19", "weather", "w"]:
            print("Heh want to check weather huh? your luck is already in thunderstorm XD")
            startweather()
        elif choice in ["20", "video", "player"]:
            video_player()
        elif choice in ["21", "exit", "quit", "q"]:
            exit_program()
            print("Bro why you oppendd the calculator just to exit it???")
            break
        else:
            print("Bro really what just you typed???")
            print("Bro im tired of you im just terminating myself now!")
            break

menu()
